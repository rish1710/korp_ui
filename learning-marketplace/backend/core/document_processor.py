"""
Document Processor — Extracts text from PDFs, PPTs, images, and YouTube videos.
"""
import os
from typing import Optional
from PyPDF2 import PdfReader
from pptx import Presentation
from PIL import Image
import pytesseract
import pdfplumber


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from a PDF file using pdfplumber (better table/layout support)."""
    text_parts = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    
    # Fallback to PyPDF2 if pdfplumber fails
    if not text_parts:
        reader = PdfReader(file_path)
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text_parts.append(t)
    
    return "\n\n".join(text_parts)


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PowerPoint files."""
    prs = Presentation(file_path)
    text_parts = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text.append(text)
        if slide_text:
            text_parts.append("\n".join(slide_text))
    return "\n\n---\n\n".join(text_parts)


def extract_text_from_image(file_path: str) -> str:
    """Extract text from images using OCR (Tesseract)."""
    image = Image.open(file_path)
    text = pytesseract.image_to_string(image)
    return text.strip()


def extract_text_from_youtube(url: str) -> str:
    """Extract transcript from a YouTube video."""
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
        
        # Extract video ID from URL
        video_id = None
        if "v=" in url:
            video_id = url.split("v=")[1].split("&")[0]
        elif "youtu.be/" in url:
            video_id = url.split("youtu.be/")[1].split("?")[0]
        
        if not video_id:
            return "Could not extract video ID from URL."
        
        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        text = " ".join([entry["text"] for entry in transcript])
        return text
    except Exception as e:
        return f"Error extracting YouTube transcript: {str(e)}"


def process_document(file_path: str, file_type: str) -> str:
    """Route document to appropriate extractor."""
    if file_type == "pdf":
        return extract_text_from_pdf(file_path)
    elif file_type in ("ppt", "pptx"):
        return extract_text_from_pptx(file_path)
    elif file_type in ("png", "jpg", "jpeg", "webp"):
        return extract_text_from_image(file_path)
    elif file_type == "youtube":
        return extract_text_from_youtube(file_path)  # file_path is URL here
    elif file_type == "txt":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    else:
        return f"Unsupported file type: {file_type}"


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    """Split text into overlapping chunks for embedding."""
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = " ".join(words[i:i + chunk_size])
        if chunk.strip():
            chunks.append(chunk)
    return chunks
